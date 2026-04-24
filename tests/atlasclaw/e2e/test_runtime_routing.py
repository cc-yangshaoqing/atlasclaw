# -*- coding: utf-8 -*-
# Copyright 2021  Qianyun, Inc. All rights reserved.

"""Real-agent E2E routing regressions with deterministic model and fake externals."""

from __future__ import annotations

import json
import re
import threading
import time
from dataclasses import dataclass
from http.server import BaseHTTPRequestHandler, ThreadingHTTPServer
from pathlib import Path
from typing import Any, Optional
from urllib.parse import parse_qs, urlparse

import pytest
from fastapi.testclient import TestClient
from pydantic_ai.messages import (
    ModelMessage,
    ModelRequest,
    ModelResponse,
    TextPart,
    ToolCallPart,
    ToolReturnPart,
    UserPromptPart,
)
from pydantic_ai.models.function import AgentInfo, DeltaToolCall, FunctionModel
from pptx import Presentation


pytestmark = pytest.mark.e2e

ATLAS_ADMIN_USER = "admin"
ATLAS_ADMIN_PASS = "Admin@123"

PENDING_APPROVALS: list[dict[str, Any]] = [
    {
        "id": "REQ-20260316-0001",
        "workflowId": "TIC20260316000001",
        "name": "扩容生产主机",
        "catalogName": "Machine Service",
        "applicant": "张三",
        "email": "zhangsan@example.com",
        "description": "Test ticket for build verification",
        "createdDate": 1773373200000,
        "updatedDate": 1773376800000,
        "chargePredictResult": {"totalCost": 1234},
        "currentActivity": {
            "id": "APR-20260316-0001",
            "taskId": "TASK-20260316-0001",
            "processInstanceId": "PROC-20260316-0001",
            "processStep": {"name": "一级审批"},
            "assignments": [{"approver": {"name": "一级审批"}}],
            "requestParams": {"resourceSpecs": {"node-1": {"cpu": 4, "memory": 8192, "disk": "100GB"}}},
        },
    },
    {
        "id": "REQ-20260313-0006",
        "workflowId": "TIC20260313000006",
        "name": "创建测试 VPC",
        "catalogName": "VPC Service",
        "applicant": "李四",
        "email": "lisi@example.com",
        "description": "Create a new VPC for QA validation",
        "createdDate": 1773110400000,
        "updatedDate": 1773114000000,
        "currentActivity": {
            "id": "APR-20260313-0006",
            "taskId": "TASK-20260313-0006",
            "processInstanceId": "PROC-20260313-0006",
            "processStep": {"name": "一级审批"},
            "assignments": [{"approver": {"name": "一级审批"}}],
            "requestParams": {"resourceSpecs": {"network": {"infra_type": "network"}}},
        },
    },
    {
        "id": "REQ-20260313-0004",
        "workflowId": "TIC20260313000004",
        "name": "开通支持服务",
        "catalogName": "Support Service",
        "applicant": "王五",
        "email": "wangwu@example.com",
        "description": "Enable support service for the project team",
        "createdDate": 1773103200000,
        "updatedDate": 1773108600000,
        "currentActivity": {
            "id": "APR-20260313-0004",
            "taskId": "TASK-20260313-0004",
            "processInstanceId": "PROC-20260313-0004",
            "processStep": {"name": "二级审批"},
            "assignments": [{"approver": {"name": "二级审批"}}],
            "requestParams": {"resourceSpecs": {"support": {"infra_type": "service"}}},
        },
    },
]

SERVICE_CATALOGS: list[dict[str, Any]] = [
    {"id": "catalog-incident", "nameZh": "Incident Ticket", "sourceKey": "request.incident.ticket", "serviceCategory": "GENERIC_SERVICE", "instructions": ""},
    {"id": "catalog-machine", "nameZh": "Machine Service", "sourceKey": "resource.iaas.machine.test_machine", "serviceCategory": "CLOUD_COMPONENT_SERVICE", "instructions": ""},
    {"id": "catalog-vpc", "nameZh": "VPC Service", "sourceKey": "resource.iaas.network.network.testvpc", "serviceCategory": "CLOUD_COMPONENT_SERVICE", "instructions": ""},
    {"id": "catalog-support", "nameZh": "Support Service", "sourceKey": "resource.support.default", "serviceCategory": "GENERIC_SERVICE", "instructions": ""},
]


@dataclass
class AgentHarness:
    """Shared objects for true-agent E2E execution."""

    client: TestClient
    headers: dict[str, str]
    workspace_path: Path
    cmp_state: dict[str, Any]

    def create_thread(self) -> str:
        response = self.client.post(
            "/api/sessions/threads",
            json={"agent_id": "main", "channel": "web", "chat_type": "dm", "account_id": "default"},
            headers=self.headers,
        )
        assert response.status_code == 200, response.text
        payload = response.json()
        session_key = str(payload.get("session_key", "")).strip()
        assert session_key, f"session_key missing: {payload}"
        return session_key


@dataclass
class AgentRunOutcome:
    """Collected artifacts from one /api/agent/run cycle."""

    run_id: str
    message: str
    session_key: str
    wall_seconds: float
    events: dict[str, list]
    assistant_text: str
    tool_starts: list[str]
    runtime_states: list[str]
    runtime_items: list[dict[str, Any]]
    final_status: str


def _build_e2e_config(tmp_path: Path, db_path: Path, cmp_base_url: str) -> dict[str, Any]:
    """Build a minimal AtlasClaw config for deterministic E2E routing tests."""
    project_root = Path(__file__).resolve().parent.parent.parent.parent
    providers_root = str((project_root.parent / "atlasclaw-providers" / "providers").resolve())
    skills_root = str((project_root.parent / "atlasclaw-providers" / "skills").resolve())
    return {
        "workspace": {"path": str((tmp_path / ".atlasclaw-e2e").resolve())},
        "providers_root": providers_root,
        "skills_root": skills_root,
        "database": {"type": "sqlite", "sqlite": {"path": str(db_path.resolve())}},
        "auth": {
            "enabled": True,
            "provider": "local",
            "jwt": {
                "secret_key": "e2e-routing-secret-key",
                "issuer": "atlasclaw-e2e-routing",
                "header_name": "AtlasClaw-Authenticate",
                "cookie_name": "AtlasClaw-Authenticate",
                "expires_minutes": 60,
            },
            "local": {
                "enabled": True,
                "default_admin_username": ATLAS_ADMIN_USER,
                "default_admin_password": ATLAS_ADMIN_PASS,
            },
        },
        "model": {
            "primary": "test-token",
            "fallbacks": [],
            "temperature": 0.0,
            "selection_strategy": "health",
            "tokens": [
                {
                    "id": "test-token",
                    "provider": "openai",
                    "model": "deterministic-e2e",
                    "base_url": "https://example.invalid/v1",
                    "api_key": "test-key-placeholder",
                    "api_type": "openai",
                    "priority": 100,
                    "weight": 100,
                }
            ],
        },
        "service_providers": {
            "smartcmp": {
                "default": {
                    "base_url": cmp_base_url,
                    "cookie": "AtlasClaw-Host-Authenticate=fake-host-token",
                }
            }
        },
    }


def _login(client: TestClient) -> dict[str, Any]:
    """Login with local auth and return the JSON body."""
    response = client.post(
        "/api/auth/local/login",
        json={"username": ATLAS_ADMIN_USER, "password": ATLAS_ADMIN_PASS},
    )
    assert response.status_code == 200, response.text
    payload = response.json()
    assert payload.get("success") is True, payload
    return payload


def _auth_headers(login_body: dict[str, Any]) -> dict[str, str]:
    """Build auth headers from login response."""
    token = str(login_body.get("token", "")).strip()
    assert token, f"login token missing: {login_body}"
    return {"AtlasClaw-Authenticate": token}


def _run_agent_and_collect_events(
    client: TestClient,
    headers: dict[str, str],
    message: str,
    session_key: str,
    timeout: int = 120,
) -> tuple[str, str, dict[str, list], float]:
    """Start a real agent run, wait for completion, then collect buffered SSE events."""
    started_at = time.perf_counter()
    run_resp = client.post(
        "/api/agent/run",
        json={"session_key": session_key, "message": message, "timeout_seconds": timeout},
        headers=headers,
    )
    assert run_resp.status_code == 200, f"Agent run failed: {run_resp.text}"
    run_data = run_resp.json()
    run_id = run_data["run_id"]
    assert run_data["status"] == "running"

    deadline = time.perf_counter() + timeout
    final_status = "running"
    while time.perf_counter() < deadline:
        time.sleep(0.1)
        status_resp = client.get(f"/api/agent/runs/{run_id}", headers=headers)
        assert status_resp.status_code == 200, status_resp.text
        status_data = status_resp.json()
        final_status = str(status_data.get("status", "unknown"))
        if final_status != "running":
            break

    events: dict[str, list] = {
        "lifecycle": [],
        "assistant": [],
        "tool": [],
        "error": [],
        "thinking": [],
        "runtime": [],
        "raw_lines": [],
    }
    with client.stream("GET", f"/api/agent/runs/{run_id}/stream", headers=headers) as stream_resp:
        assert stream_resp.status_code == 200, stream_resp.text
        current_event_type = ""
        current_data = ""
        for line in stream_resp.iter_lines():
            events["raw_lines"].append(line)
            if line.startswith("event:"):
                current_event_type = line[len("event:") :].strip()
                continue
            if line.startswith("data:"):
                current_data = line[len("data:") :].strip()
                continue
            if line != "" or not current_event_type or not current_data:
                continue
            try:
                payload = json.loads(current_data)
            except json.JSONDecodeError:
                payload = {"raw": current_data}
            if current_event_type in events:
                events[current_event_type].append(payload)
            if current_event_type == "lifecycle" and payload.get("phase") in {"end", "error", "timeout"}:
                break
            current_event_type = ""
            current_data = ""

    events["_final_status"] = [final_status]
    wall_seconds = round(time.perf_counter() - started_at, 3)
    return run_id, final_status, events, wall_seconds


def _get_assistant_text(events: dict[str, list]) -> str:
    """Concatenate streamed assistant chunks."""
    return "".join(str(evt.get("text", "")) for evt in events.get("assistant", []) if evt.get("text"))


def _runtime_states(events: dict[str, list]) -> list[str]:
    """Extract runtime state names in stream order."""
    return [str(evt.get("state", "")).strip() for evt in events.get("runtime", []) if str(evt.get("state", "")).strip()]


def _tool_starts(events: dict[str, list]) -> list[str]:
    """Return tool names for tool:start events in stream order."""
    tool_names: list[str] = []
    for evt in events.get("tool", []):
        if evt.get("phase") == "start":
            name = str(evt.get("tool", "")).strip()
            if name:
                tool_names.append(name)
    return tool_names


def _assert_completed(outcome: AgentRunOutcome) -> None:
    """Assert the agent completed successfully and emitted an assistant answer."""
    lifecycle = [evt.get("phase") for evt in outcome.events.get("lifecycle", [])]
    errors = outcome.events.get("error", [])
    assert not errors, f"stream errors: {errors}"
    assert lifecycle[:1] == ["start"], lifecycle
    assert lifecycle[-1:] == ["end"], lifecycle
    assert outcome.final_status == "completed", outcome.final_status
    assert outcome.runtime_states[-1:] == ["answered"], outcome.runtime_states
    assert outcome.assistant_text.strip(), f"assistant text missing for {outcome.run_id}"


def _extract_json_block(text: str, start_marker: str, end_marker: str) -> Any:
    """Extract JSON payload between tool output markers."""
    pattern = re.escape(start_marker) + r"\s*(.*?)\s*" + re.escape(end_marker)
    match = re.search(pattern, text, flags=re.DOTALL)
    if not match:
        return None
    return json.loads(match.group(1))


def _tool_output_text(tool_content: Any) -> str:
    """Extract the human-readable output string from a tool return payload."""
    if isinstance(tool_content, dict):
        output = tool_content.get("output")
        if isinstance(output, str):
            return output
        content_blocks = tool_content.get("content")
        if isinstance(content_blocks, list):
            texts = [str(block.get("text", "")) for block in content_blocks if isinstance(block, dict) and block.get("type") == "text"]
            return "\n".join(texts).strip()
    if isinstance(tool_content, str):
        return tool_content
    return ""


def _tool_internal_payload(tool_content: Any) -> Any:
    """Extract structured internal metadata from a tool return payload."""
    if not isinstance(tool_content, dict):
        return None
    internal = tool_content.get("_internal")
    if isinstance(internal, (dict, list)):
        return internal
    if isinstance(internal, str) and internal.strip():
        return json.loads(internal)
    return None


def _extract_catalog_meta(content: Any) -> list[dict[str, Any]]:
    """Normalize catalog metadata from either legacy stdout or current _internal payloads."""
    internal = _tool_internal_payload(content)
    if isinstance(internal, list):
        return [item for item in internal if isinstance(item, dict)]
    if isinstance(internal, dict):
        catalogs = internal.get("catalogs")
        if isinstance(catalogs, list):
            return [item for item in catalogs if isinstance(item, dict)]

    output = _tool_output_text(content)
    parsed = _extract_json_block(output, "##CATALOG_META_START##", "##CATALOG_META_END##")
    if isinstance(parsed, list):
        return [item for item in parsed if isinstance(item, dict)]
    if isinstance(parsed, dict):
        catalogs = parsed.get("catalogs")
        if isinstance(catalogs, list):
            return [item for item in catalogs if isinstance(item, dict)]
    return []


def _iter_message_parts(messages: list[ModelMessage]) -> list[Any]:
    """Flatten model message parts for simpler history inspection."""
    parts: list[Any] = []
    for message in messages:
        parts.extend(list(getattr(message, "parts", []) or []))
    return parts


def _latest_model_request(messages: list[ModelMessage]) -> Optional[ModelRequest]:
    """Return the most recent ModelRequest from message history."""
    for message in reversed(messages):
        if isinstance(message, ModelRequest):
            return message
    return None


def _extract_latest_user_text(messages: list[ModelMessage]) -> str:
    """Read the user prompt from the most recent ModelRequest only."""
    latest_request = _latest_model_request(messages)
    if latest_request is None:
        return ""
    for part in reversed(list(latest_request.parts)):
        if isinstance(part, UserPromptPart) and isinstance(part.content, str):
            return part.content.strip()
    return ""


def _extract_previous_assistant_text(messages: list[ModelMessage]) -> str:
    """Read the latest assistant text already present in model history."""
    fallback_chunks: list[str] = []
    for message in reversed(messages):
        if not isinstance(message, ModelResponse):
            for part in getattr(message, "parts", []) or []:
                if isinstance(part, TextPart) and isinstance(part.content, str) and part.content.strip():
                    fallback_chunks.append(part.content.strip())
            continue
        combined = "".join(part.content for part in message.parts if isinstance(part, TextPart)).strip()
        if combined:
            return combined
    return "\n".join(reversed(fallback_chunks)).strip()


def _history_contains_weather_context(messages: list[ModelMessage]) -> bool:
    """Return whether recent conversation history already contains a weather exchange."""
    history_text = "\n".join(
        chunk
        for chunk in [
            _extract_previous_assistant_text(messages),
            "\n".join(
                part.content.strip()
                for part in _iter_message_parts(messages)
                if isinstance(part, UserPromptPart)
                and isinstance(part.content, str)
                and part.content.strip()
            ),
        ]
        if chunk
    ).strip()
    if not history_text:
        return False
    lowered = history_text.lower()
    return "weather for" in lowered or "天气" in history_text or "forecast" in lowered


def _extract_latest_tool_return(messages: list[ModelMessage]) -> Optional[ToolReturnPart]:
    """Find a tool return on the most recent ModelRequest only."""
    latest_request = _latest_model_request(messages)
    if latest_request is None:
        return None
    for part in reversed(list(latest_request.parts)):
        if isinstance(part, ToolReturnPart):
            return part
    return None


def _extract_pending_items(messages: list[ModelMessage]) -> list[dict[str, Any]]:
    """Recover pending request items from prior tool results or assistant history."""
    for part in reversed(_iter_message_parts(messages)):
        if isinstance(part, ToolReturnPart) and part.tool_name == "smartcmp_list_pending":
            output = _tool_output_text(part.content)
            parsed = _extract_json_block(output, "##APPROVAL_META_START##", "##APPROVAL_META_END##")
            if isinstance(parsed, list):
                return [item for item in parsed if isinstance(item, dict)]

    assistant_text = _extract_previous_assistant_text(messages)
    raw_history_chunks = [assistant_text] if assistant_text else []
    for message in messages:
        model_dump_json = getattr(message, "model_dump_json", None)
        if callable(model_dump_json):
            try:
                raw_history_chunks.append(str(model_dump_json()))
                continue
            except Exception:
                pass
        raw_history_chunks.append(repr(message))
    history_corpus = "\n".join(chunk for chunk in raw_history_chunks if chunk).strip()
    matched: list[dict[str, Any]] = []
    for item in PENDING_APPROVALS:
        workflow_id = str(item.get("workflowId", "")).strip()
        if workflow_id and workflow_id in history_corpus:
            matched.append(
                {
                    "id": workflow_id,
                    "title": item.get("name", ""),
                    "approver": ((((item.get("currentActivity") or {}).get("assignments") or [{}])[0]).get("approver", {}).get("name", "")),
                    "approvalId": (item.get("currentActivity") or {}).get("id", ""),
                    "summary": item.get("description", ""),
                }
            )
    return matched


def _require_tool(available_tools: set[str], tool_name: str, user_text: str) -> None:
    """Fail fast when routing does not surface the expected tool."""
    assert tool_name in available_tools, (
        f"Expected tool {tool_name!r} to be available for {user_text!r}, "
        f"got {sorted(available_tools)!r}"
    )


def _format_pending_summary(content: Any) -> str:
    output = _tool_output_text(content)
    parsed = _extract_json_block(output, "##APPROVAL_META_START##", "##APPROVAL_META_END##")
    assert isinstance(parsed, list) and parsed, f"pending meta missing: {content!r}"
    items = []
    for item in parsed[:3]:
        workflow_id = str(item.get("workflowId", "")).strip()
        name = str(item.get("name", "")).strip()
        approver = str(item.get("currentApprover", "")).strip()
        fragment = f"{workflow_id} {name}"
        if approver:
            fragment += f"（{approver}）"
        items.append(fragment)
    return f"当前待审批申请共 {len(parsed)} 条：" + "，".join(items) + "。"


def _format_detail_summary(content: Any) -> str:
    output = _tool_output_text(content)
    parsed = _extract_json_block(output, "##APPROVAL_DETAIL_META_START##", "##APPROVAL_DETAIL_META_END##")
    assert isinstance(parsed, dict) and parsed, f"detail meta missing: {content!r}"
    return (
        f"{parsed.get('workflowId', '')} 的详情：{parsed.get('description', '')}，"
        f"标题是 {parsed.get('name', '')}，Approval ID 是 {parsed.get('approvalId', '')}。"
    )


def _format_services_summary(content: Any) -> str:
    parsed = _extract_catalog_meta(content)
    assert parsed, f"catalog meta missing: {content!r}"
    names = [str(item.get("name", "")).strip() for item in parsed if str(item.get("name", "")).strip()]
    return "当前服务目录包括：" + "、".join(names) + "。"


def _format_weather_summary(content: Any) -> str:
    assert isinstance(content, dict), content
    markdown = ""
    for block in content.get("content", []) or []:
        if isinstance(block, dict) and block.get("type") == "text":
            markdown = str(block.get("text", ""))
            break
    assert markdown, f"weather markdown missing: {content!r}"
    line = next((row for row in markdown.splitlines() if row.startswith("| 2026-04-13 |")), "")
    assert line, markdown
    columns = [part.strip() for part in line.split("|") if part.strip()]
    location = ((content.get("details") or {}).get("resolved_location") or {}).get("name") or "上海"
    return f"{location} 今天天气晴到多云，最低 {columns[2]}，最高 {columns[3]}，降水 {columns[4]}，风速 {columns[6]}。"


def _format_ppt_summary(content: Any) -> str:
    assert isinstance(content, dict) and content.get("success") is True, content
    return f"已经根据上文申请生成新的 PPT 文件，共 {content.get('slide_count')} 页：{content.get('file_path')}"


def _direct_park_answer() -> str:
    return (
        "上海周边适合休闲骑行的公园和绿道，可以先看这几个方向：浦东滨江骑行道、"
        "世纪公园周边绿道、青西郊野公园、共青森林公园附近骑行线，以及苏州河健身步道一带。"
    )


def _decide_model_action(messages: list[ModelMessage], agent_info: AgentInfo) -> tuple[str, Any]:
    """Return either ('text', text) or ('tool', name, args, call_id)."""
    available_tools = {tool.name for tool in agent_info.function_tools}
    last_tool_return = _extract_latest_tool_return(messages)
    if last_tool_return is not None:
        if last_tool_return.tool_name == "smartcmp_list_pending":
            return "text", _format_pending_summary(last_tool_return.content)
        if last_tool_return.tool_name == "smartcmp_get_request_detail":
            return "text", _format_detail_summary(last_tool_return.content)
        if last_tool_return.tool_name == "smartcmp_list_services":
            return "text", _format_services_summary(last_tool_return.content)
        if last_tool_return.tool_name == "openmeteo_weather":
            return "text", _format_weather_summary(last_tool_return.content)
        if last_tool_return.tool_name == "pptx_create_deck":
            return "text", _format_ppt_summary(last_tool_return.content)
        raise AssertionError(f"Unexpected tool continuation: {last_tool_return.tool_name}")

    user_text = _extract_latest_user_text(messages)
    if "待审批" in user_text or "审批数据" in user_text:
        _require_tool(available_tools, "smartcmp_list_pending", user_text)
        return "tool", "smartcmp_list_pending", {}, "call-pending-1"
    if "详情" in user_text and "TIC20260316000001" in user_text:
        _require_tool(available_tools, "smartcmp_get_request_detail", user_text)
        return "tool", "smartcmp_get_request_detail", {"identifier": "TIC20260316000001"}, "call-detail-1"
    if "服务目录" in user_text:
        _require_tool(available_tools, "smartcmp_list_services", user_text)
        return "tool", "smartcmp_list_services", {}, "call-services-1"
    if "天气" in user_text:
        _require_tool(available_tools, "openmeteo_weather", user_text)
        return "tool", "openmeteo_weather", {"location": "上海", "days": 2}, "call-weather-1"
    if user_text == "上海呢" and _history_contains_weather_context(messages):
        _require_tool(available_tools, "openmeteo_weather", user_text)
        return "tool", "openmeteo_weather", {"location": "上海", "days": 2}, "call-weather-follow-up-1"
    if "骑行公园" in user_text:
        return "text", _direct_park_answer()
    if "PPT" in user_text.upper():
        _require_tool(available_tools, "pptx_create_deck", user_text)
        items = _extract_pending_items(messages)
        assert items, "Follow-up PPT request did not retain prior pending request context."
        return (
            "tool",
            "pptx_create_deck",
            {
                "title": "CMP 待审批申请汇总",
                "subtitle": "由真实 Agent E2E 测试生成",
                "items": items,
                "output_filename": "pending-approvals-e2e.pptx",
            },
            "call-pptx-1",
        )
    return "text", f"Unhandled deterministic E2E prompt: {user_text}"


def _function_response(messages: list[ModelMessage], agent_info: AgentInfo) -> ModelResponse:
    """Non-streamed deterministic model response used by agent.iter()."""
    decision = _decide_model_action(messages, agent_info)
    if decision[0] == "tool":
        _, tool_name, args, call_id = decision
        return ModelResponse(parts=[ToolCallPart(tool_name, args, tool_call_id=call_id)])
    return ModelResponse(parts=[TextPart(decision[1])])


async def _stream_function_response(messages: list[ModelMessage], agent_info: AgentInfo):
    """Streamed deterministic response for any code path that requests streaming."""
    decision = _decide_model_action(messages, agent_info)
    if decision[0] == "tool":
        _, tool_name, args, call_id = decision
        yield {0: DeltaToolCall(name=tool_name, json_args=json.dumps(args, ensure_ascii=False), tool_call_id=call_id)}
        return
    yield decision[1]


def _create_test_model(_token: Any) -> FunctionModel:
    """Factory patched into create_pydantic_model for true-agent E2E tests."""
    return FunctionModel(function=_function_response, stream_function=_stream_function_response, model_name="deterministic-routing-e2e")


async def _fake_weather_request_json(*, url: str, params: dict[str, Any], timeout_seconds: float = 12.0) -> dict[str, Any]:
    """Fake Open-Meteo responses while keeping the real weather tool logic intact."""
    _ = (params, timeout_seconds)
    if "geocoding-api.open-meteo.com" in url:
        return {"results": [{"name": "Shanghai", "country": "China", "admin1": "Shanghai", "latitude": 31.2304, "longitude": 121.4737}]}
    if "api.open-meteo.com" in url:
        return {
            "current": {"time": "2026-04-13T09:00", "temperature_2m": 22.0, "wind_speed_10m": 12.0, "weather_code": 1},
            "current_units": {"temperature_2m": "°C", "wind_speed_10m": "kmh"},
            "daily": {
                "time": ["2026-04-13", "2026-04-14"],
                "temperature_2m_max": [26.0, 27.0],
                "temperature_2m_min": [18.0, 19.0],
                "weather_code": [1, 2],
                "precipitation_sum": [0.0, 0.2],
                "precipitation_probability_max": [10, 20],
                "wind_speed_10m_max": [18.0, 20.0],
            },
            "daily_units": {"temperature_2m_max": "°C", "temperature_2m_min": "°C", "precipitation_sum": "mm", "precipitation_probability_max": "%", "wind_speed_10m_max": "kmh"},
        }
    raise AssertionError(f"Unexpected weather URL: {url}")


class _FakeCmpHandler(BaseHTTPRequestHandler):
    """Serve minimal SmartCMP JSON payloads expected by provider scripts."""

    server_version = "FakeSmartCMP/1.0"

    def do_GET(self) -> None:  # noqa: N802
        parsed = urlparse(self.path)
        query = parse_qs(parsed.query)
        state = getattr(self.server, "state")
        state["requests"].append({"path": parsed.path, "query": query})
        if parsed.path.endswith("/generic-request/current-activity-approval"):
            state["approvals_requests"] += 1
            self._write_json({"content": PENDING_APPROVALS, "totalElements": len(PENDING_APPROVALS)})
            return
        if parsed.path.endswith("/catalogs/published"):
            state["catalog_requests"] += 1
            self._write_json({"content": SERVICE_CATALOGS, "totalElements": len(SERVICE_CATALOGS)})
            return
        self.send_response(404)
        self.send_header("Content-Type", "application/json; charset=utf-8")
        self.end_headers()
        self.wfile.write(b'{"error":"not-found"}')

    def log_message(self, format: str, *args: Any) -> None:  # noqa: A003
        return

    def _write_json(self, payload: dict[str, Any]) -> None:
        body = json.dumps(payload, ensure_ascii=False).encode("utf-8")
        self.send_response(200)
        self.send_header("Content-Type", "application/json; charset=utf-8")
        self.send_header("Content-Length", str(len(body)))
        self.end_headers()
        self.wfile.write(body)


@pytest.fixture()
def fake_cmp_server():
    """Start a local fake SmartCMP server for one test."""
    state = {"approvals_requests": 0, "catalog_requests": 0, "requests": []}
    server = ThreadingHTTPServer(("127.0.0.1", 0), _FakeCmpHandler)
    server.state = state
    thread = threading.Thread(target=server.serve_forever, daemon=True)
    thread.start()
    base_url = f"http://127.0.0.1:{server.server_address[1]}"
    try:
        yield {"base_url": base_url, "state": state}
    finally:
        server.shutdown()
        thread.join(timeout=5)
        server.server_close()


@pytest.fixture()
def agent_harness(tmp_path: Path, monkeypatch: pytest.MonkeyPatch, fake_cmp_server) -> AgentHarness:
    """Create the real app with deterministic model and fake external services."""
    db_path = tmp_path / "runtime-routing.db"
    config_path = tmp_path / "atlasclaw.runtime-routing-e2e.json"
    config = _build_e2e_config(tmp_path, db_path, fake_cmp_server["base_url"])
    config_path.write_text(json.dumps(config, ensure_ascii=False, indent=2), encoding="utf-8")
    monkeypatch.setenv("ATLASCLAW_CONFIG", str(config_path.resolve()))

    import app.atlasclaw.bootstrap.startup_helpers as startup_helpers
    import app.atlasclaw.core.config as config_module
    import app.atlasclaw.main as main_module
    import app.atlasclaw.tools.web.openmeteo_weather_tool as weather_module

    old_config_manager = config_module._config_manager
    config_module._config_manager = config_module.ConfigManager(config_path=str(config_path.resolve()))
    monkeypatch.setattr(main_module, "create_pydantic_model", _create_test_model)
    monkeypatch.setattr(startup_helpers, "create_pydantic_model", _create_test_model)
    monkeypatch.setattr(weather_module, "_request_json", _fake_weather_request_json)

    app = main_module.create_app()
    try:
        with TestClient(app) as client:
            login_body = _login(client)
            headers = _auth_headers(login_body)
            yield AgentHarness(
                client=client,
                headers=headers,
                workspace_path=Path(config["workspace"]["path"]).resolve(),
                cmp_state=fake_cmp_server["state"],
            )
    finally:
        config_module._config_manager = old_config_manager


def _run_round(harness: AgentHarness, message: str, session_key: str) -> AgentRunOutcome:
    """Execute one real agent round and normalize the collected outputs."""
    run_id, final_status, events, wall_seconds = _run_agent_and_collect_events(
        harness.client,
        harness.headers,
        message,
        session_key,
    )
    return AgentRunOutcome(
        run_id=run_id,
        message=message,
        session_key=session_key,
        wall_seconds=wall_seconds,
        events=events,
        assistant_text=_get_assistant_text(events),
        tool_starts=_tool_starts(events),
        runtime_states=_runtime_states(events),
        runtime_items=events.get("runtime", []),
        final_status=final_status,
    )


def test_cmp_pending_query_runs_real_agent_loop(agent_harness: AgentHarness) -> None:
    session_key = agent_harness.create_thread()
    outcome = _run_round(agent_harness, "查一个 cmp 所有待审批的申请", session_key)

    _assert_completed(outcome)
    assert outcome.tool_starts == ["smartcmp_list_pending"]
    assert "TIC20260316000001" in outcome.assistant_text
    assert "TIC20260313000006" in outcome.assistant_text
    assert "TIC20260313000004" in outcome.assistant_text
    assert "一级审批" in outcome.assistant_text
    assert agent_harness.cmp_state["approvals_requests"] >= 1


def test_cmp_detail_query_runs_real_agent_loop(agent_harness: AgentHarness) -> None:
    session_key = agent_harness.create_thread()
    outcome = _run_round(agent_harness, "我要看下TIC20260316000001的详情", session_key)

    _assert_completed(outcome)
    assert outcome.tool_starts == ["smartcmp_get_request_detail"]
    assert "TIC20260316000001" in outcome.assistant_text
    assert "Test ticket for build verification" in outcome.assistant_text
    assert "APR-20260316-0001" in outcome.assistant_text
    assert agent_harness.cmp_state["approvals_requests"] >= 1


def test_cmp_services_query_runs_real_agent_loop(agent_harness: AgentHarness) -> None:
    session_key = agent_harness.create_thread()
    outcome = _run_round(agent_harness, "还有查下CMP 里目前有的服务目录", session_key)

    _assert_completed(outcome)
    assert outcome.tool_starts == ["smartcmp_list_services"]
    assert "Incident Ticket" in outcome.assistant_text
    assert "Machine Service" in outcome.assistant_text
    assert "VPC Service" in outcome.assistant_text
    assert "Support Service" in outcome.assistant_text
    assert agent_harness.cmp_state["catalog_requests"] >= 1


def test_weather_query_runs_real_agent_loop(agent_harness: AgentHarness) -> None:
    session_key = agent_harness.create_thread()
    outcome = _run_round(agent_harness, "帮我查下上海天气", session_key)

    _assert_completed(outcome)
    assert outcome.tool_starts == ["openmeteo_weather"]
    assert "Shanghai" in outcome.assistant_text or "上海" in outcome.assistant_text
    assert "26.0°C" in outcome.assistant_text
    assert "18.0kmh" in outcome.assistant_text
    assert outcome.runtime_states.count("reasoning") >= 2


def test_weather_follow_up_short_location_reuses_context_and_runs_weather_tool(
    agent_harness: AgentHarness,
) -> None:
    session_key = agent_harness.create_thread()
    first = _run_round(agent_harness, "明天北京天气呢", session_key)
    _assert_completed(first)
    assert first.tool_starts == ["openmeteo_weather"]

    second = _run_round(agent_harness, "上海呢", session_key)
    _assert_completed(second)
    assert second.tool_starts == ["openmeteo_weather"]
    assert "Shanghai" in second.assistant_text or "上海" in second.assistant_text


def test_public_park_query_answers_directly_without_tool(agent_harness: AgentHarness) -> None:
    session_key = agent_harness.create_thread()
    outcome = _run_round(agent_harness, "我想查下上海周边的骑行公园", session_key)

    _assert_completed(outcome)
    assert outcome.tool_starts == []
    assert "世纪公园" in outcome.assistant_text
    assert "浦东滨江" in outcome.assistant_text
    assert "青西郊野公园" in outcome.assistant_text


def test_pending_then_ppt_follow_up_creates_real_pptx(agent_harness: AgentHarness) -> None:
    session_key = agent_harness.create_thread()
    first = _run_round(agent_harness, "查一个 cmp 所有待审批的申请", session_key)
    _assert_completed(first)
    approvals_before_follow_up = agent_harness.cmp_state["approvals_requests"]

    second = _run_round(agent_harness, "将这些申请写入一个新的PPT", session_key)
    _assert_completed(second)

    assert second.tool_starts == ["pptx_create_deck"]
    assert agent_harness.cmp_state["approvals_requests"] == approvals_before_follow_up
    assert "pending-approvals-e2e.pptx" in second.assistant_text

    pptx_match = re.search(r"([A-Za-z]:\\\S+\.pptx|/\S+\.pptx)", second.assistant_text)
    assert pptx_match, second.assistant_text
    pptx_path = Path(pptx_match.group(1))
    assert pptx_path.is_file(), f"PPTX file not created: {pptx_path}"
    assert pptx_path.suffix.lower() == ".pptx"
    assert str(agent_harness.workspace_path) in str(pptx_path)

    presentation = Presentation(str(pptx_path))
    assert len(presentation.slides) == 5
    title_shapes = [shape.text for shape in presentation.slides[0].shapes if hasattr(shape, "text")]
    assert any("CMP 待审批申请汇总" in text for text in title_shapes)


def test_pending_then_english_ppt_follow_up_creates_real_pptx(agent_harness: AgentHarness) -> None:
    session_key = agent_harness.create_thread()
    first = _run_round(agent_harness, "查下CMP现在的审批数据", session_key)
    _assert_completed(first)
    assert first.tool_starts == ["smartcmp_list_pending"]
    approvals_before_follow_up = agent_harness.cmp_state["approvals_requests"]
    assert approvals_before_follow_up >= 1

    second = _run_round(agent_harness, "write the request data into a PPT", session_key)
    _assert_completed(second)

    assert second.tool_starts == ["pptx_create_deck"]
    assert agent_harness.cmp_state["approvals_requests"] == approvals_before_follow_up
    assert ".pptx" in second.assistant_text.lower()
    assert ".txt" not in second.assistant_text.lower()

    pptx_match = re.search(r"([A-Za-z]:\\\S+\.pptx|/\S+\.pptx)", second.assistant_text)
    assert pptx_match, second.assistant_text
    pptx_path = Path(pptx_match.group(1))
    assert pptx_path.is_file(), f"PPTX file not created: {pptx_path}"

    presentation = Presentation(str(pptx_path))
    assert len(presentation.slides) == 5
