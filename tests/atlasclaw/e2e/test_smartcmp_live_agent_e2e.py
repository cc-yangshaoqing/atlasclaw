# -*- coding: utf-8 -*-
# Copyright 2021  Qianyun, Inc. All rights reserved.

from __future__ import annotations

import json
import os
import re
import tempfile
import time
from dataclasses import dataclass
from pathlib import Path
from typing import Any
from urllib.parse import urlparse

import pytest
import requests
from pptx import Presentation


DEFAULT_BASE_URL = "http://127.0.0.1:8000"
DEFAULT_USERNAME = "admin"
DEFAULT_PASSWORD = "admin"
DEFAULT_REMOTE_CONTAINER_NAME = "atlasclaw"


@dataclass
class LiveRunResult:
    run_id: str
    message: str
    wall_seconds: float
    lifecycle: list[str]
    runtime_states: list[str]
    runtime_items: list[dict[str, Any]]
    errors: list[str]
    assistant_text: str


def _require_live_e2e() -> tuple[str, str, str]:
    if os.getenv("ATLASCLAW_LIVE_E2E", "").strip() != "1":
        pytest.skip("Set ATLASCLAW_LIVE_E2E=1 to run SmartCMP live agent E2E.")
    return (
        os.getenv("ATLASCLAW_LIVE_BASE_URL", DEFAULT_BASE_URL).strip() or DEFAULT_BASE_URL,
        os.getenv("ATLASCLAW_LIVE_USERNAME", DEFAULT_USERNAME).strip() or DEFAULT_USERNAME,
        os.getenv("ATLASCLAW_LIVE_PASSWORD", DEFAULT_PASSWORD).strip() or DEFAULT_PASSWORD,
    )


def _parse_sse(response: requests.Response):
    event_name: str | None = None
    data_lines: list[str] = []
    for raw_line in response.iter_lines(decode_unicode=True):
        if raw_line is None:
            continue
        line = raw_line.strip("\n")
        if not line:
            if event_name is not None and data_lines:
                yield event_name, "\n".join(data_lines)
            event_name = None
            data_lines = []
            continue
        if line.startswith("event:"):
            event_name = line.split(":", 1)[1].strip()
        elif line.startswith("data:"):
            data_lines.append(line.split(":", 1)[1].strip())


def _login(session: requests.Session, *, base_url: str, username: str, password: str) -> str:
    response = session.post(
        f"{base_url}/api/auth/local/login",
        json={"username": username, "password": password},
        timeout=20,
    )
    response.raise_for_status()
    payload = response.json()
    token = str(payload.get("token", "")).strip()
    assert token, f"login token missing: {payload}"
    return token


def _create_thread(session: requests.Session, *, base_url: str, token: str) -> str:
    response = session.post(
        f"{base_url}/api/sessions/threads",
        headers={"AtlasClaw-Authenticate": token},
        json={"agent_id": "main", "channel": "web", "chat_type": "dm", "account_id": "default"},
        timeout=20,
    )
    response.raise_for_status()
    payload = response.json()
    session_key = str(payload.get("session_key", "")).strip()
    assert session_key, f"session_key missing: {payload}"
    return session_key


def _login_and_create_thread(
    *,
    base_url: str,
    username: str,
    password: str,
) -> tuple[requests.Session, str, str]:
    session = requests.Session()
    token = _login(session, base_url=base_url, username=username, password=password)
    session_key = _create_thread(session, base_url=base_url, token=token)
    return session, token, session_key


def _run_round(
    session: requests.Session,
    *,
    base_url: str,
    token: str,
    session_key: str,
    message: str,
) -> LiveRunResult:
    headers = {"AtlasClaw-Authenticate": token}
    started_at = time.perf_counter()

    response = session.post(
        f"{base_url}/api/agent/run",
        headers=headers,
        json={"session_key": session_key, "message": message, "timeout_seconds": 120},
        timeout=30,
    )
    response.raise_for_status()
    run_id = response.json()["run_id"]

    assistant_chunks: list[str] = []
    runtime_states: list[str] = []
    runtime_items: list[dict[str, Any]] = []
    lifecycle: list[str] = []
    errors: list[str] = []

    with session.get(
        f"{base_url}/api/agent/runs/{run_id}/stream",
        headers=headers,
        stream=True,
        timeout=180,
    ) as stream_response:
        stream_response.raise_for_status()
        for event_name, payload in _parse_sse(stream_response):
            try:
                obj = json.loads(payload)
            except Exception:
                obj = {"raw": payload}
            if event_name == "assistant":
                text = str(obj.get("text", ""))
                if text:
                    assistant_chunks.append(text)
            elif event_name == "runtime":
                state = str(obj.get("state", "")).strip()
                if state:
                    runtime_states.append(state)
                runtime_items.append(obj)
            elif event_name == "lifecycle":
                phase = str(obj.get("phase", "")).strip()
                if phase:
                    lifecycle.append(phase)
                if phase == "end":
                    break
            elif event_name == "error":
                errors.append(str(obj.get("message", obj)))

    return LiveRunResult(
        run_id=run_id,
        message=message,
        wall_seconds=round(time.perf_counter() - started_at, 3),
        lifecycle=lifecycle,
        runtime_states=runtime_states,
        runtime_items=runtime_items,
        errors=errors,
        assistant_text="".join(assistant_chunks).strip(),
    )


def _looks_like_transient_provider_failure(result: LiveRunResult) -> bool:
    combined = "\n".join(
        [
            result.assistant_text,
            *result.errors,
            *[str(item.get("message", "")) for item in result.runtime_items],
        ]
    ).lower()
    return (
        "connectionreseterror" in combined
        or "远程主机强迫关闭了一个现有的连接" in combined
        or "forcibly closed by the remote host" in combined
    )


def _run_round_with_transient_retry(
    session: requests.Session,
    *,
    base_url: str,
    token: str,
    session_key: str,
    message: str,
    retries: int = 2,
) -> LiveRunResult:
    attempts = max(1, retries + 1)
    last_result: LiveRunResult | None = None
    for attempt in range(1, attempts + 1):
        result = _run_round(
            session,
            base_url=base_url,
            token=token,
            session_key=session_key,
            message=message,
        )
        last_result = result
        if not _looks_like_transient_provider_failure(result):
            return result
        if attempt < attempts:
            time.sleep(1.5)
    assert last_result is not None
    return last_result


def _assert_elapsed_monotonic(result: LiveRunResult) -> None:
    last_elapsed = -1.0
    for item in result.runtime_items:
        elapsed = float(item.get("elapsed", 0.0) or 0.0)
        phase = str(item.get("phase", "") or "").strip().lower()
        if phase == "start" and elapsed <= 0.1:
            last_elapsed = -1.0
        assert elapsed >= last_elapsed, f"elapsed not monotonic in {result.run_id}: {result.runtime_items}"
        last_elapsed = elapsed


def _assert_no_timeout_or_failure(result: LiveRunResult) -> None:
    assert "failed" not in result.runtime_states, f"run failed: {result}"
    assert not result.errors, f"stream errors: {result.errors}"
    for item in result.runtime_items:
        message = str(item.get("message", "")).lower()
        assert "timeout" not in message, f"timeout detected: {item}"


def _assert_real_tool_call(result: LiveRunResult, expected_tool_name: str) -> None:
    planned_tools: list[str] = []
    for item in result.runtime_items:
        tools = item.get("tools")
        if isinstance(tools, list):
            planned_tools.extend(str(tool).strip() for tool in tools if str(tool).strip())
    assert expected_tool_name in planned_tools, f"expected tool {expected_tool_name}, got {planned_tools}"


def _assert_common_runtime_shape(result: LiveRunResult, *, require_tool: bool = True) -> None:
    assert result.lifecycle[:1] == ["start"], f"missing lifecycle start: {result.lifecycle}"
    assert result.lifecycle[-1:] == ["end"], f"missing lifecycle end: {result.lifecycle}"
    assert result.runtime_states[-1:] == ["answered"], f"final state not answered: {result.runtime_states}"
    if require_tool:
        assert "waiting_for_tool" in result.runtime_states, f"no tool phase: {result.runtime_states}"
    else:
        assert "waiting_for_tool" not in result.runtime_states, f"unexpected tool phase: {result.runtime_states}"
    assert result.assistant_text, f"assistant text missing for run {result.run_id}"
    _assert_elapsed_monotonic(result)
    _assert_no_timeout_or_failure(result)


def _extract_cmp_request_ids(text: str) -> list[str]:
    seen: set[str] = set()
    request_ids: list[str] = []
    for match in re.finditer(r"\bTIC\d{8,}\b", text or "", flags=re.IGNORECASE):
        request_id = str(match.group(0) or "").strip()
        if not request_id or request_id in seen:
            continue
        seen.add(request_id)
        request_ids.append(request_id)
    return request_ids


def _extract_generated_pptx_reference(text: str) -> tuple[str, Path]:
    match = re.search(
        r"([A-Za-z]:[\\/][^`\r\n]+?\.pptx|\.atlasclaw[\\/][^`\r\n]+?\.pptx|/[^`\r\n]+?\.pptx)",
        text,
    )
    if match:
        raw_path = match.group(1)
        if raw_path.startswith(".atlasclaw/") or raw_path.startswith(".atlasclaw\\"):
            normalized_relative = raw_path.replace("\\", "/")
            resolved = Path("C:/Projects/cmps/atlasclaw") / normalized_relative
            return raw_path, resolved
        return raw_path, Path(raw_path)

    filename_match = re.search(r"([^\\/\s`]+\.pptx)", text, flags=re.IGNORECASE)
    assert filename_match, text
    resolved = Path("C:/Projects/cmps/atlasclaw/.atlasclaw/users/admin/exports") / filename_match.group(1)
    return filename_match.group(1), resolved


def _download_remote_container_file(
    *,
    base_url: str,
    container_path: Path,
) -> Path:
    ssh_user = os.getenv("ATLASCLAW_LIVE_SSH_USERNAME", "").strip()
    ssh_password = os.getenv("ATLASCLAW_LIVE_SSH_PASSWORD", "").strip()
    container_name = (
        os.getenv("ATLASCLAW_LIVE_CONTAINER_NAME", DEFAULT_REMOTE_CONTAINER_NAME).strip()
        or DEFAULT_REMOTE_CONTAINER_NAME
    )
    assert ssh_user and ssh_password, (
        "Set ATLASCLAW_LIVE_SSH_USERNAME and ATLASCLAW_LIVE_SSH_PASSWORD "
        "to validate generated files on a remote live host."
    )

    parsed = urlparse(base_url)
    ssh_host = os.getenv("ATLASCLAW_LIVE_SSH_HOST", "").strip() or parsed.hostname or ""
    assert ssh_host, f"Could not resolve SSH host from base_url={base_url!r}"

    import base64
    import importlib

    paramiko = importlib.import_module("paramiko")
    command = (
        "docker exec "
        f"{container_name} "
        "python -c "
        f"\"import base64, pathlib; p = pathlib.Path({container_path.as_posix()!r}); "
        "assert p.is_file(), str(p); "
        "print(base64.b64encode(p.read_bytes()).decode())\""
    )

    client = paramiko.SSHClient()
    client.set_missing_host_key_policy(paramiko.AutoAddPolicy())
    client.connect(ssh_host, username=ssh_user, password=ssh_password, timeout=20)
    try:
        _, stdout, stderr = client.exec_command(command, timeout=180)
        encoded = stdout.read().decode("utf-8", "ignore").strip()
        err = stderr.read().decode("utf-8", "ignore").strip()
    finally:
        client.close()

    assert encoded, f"remote file fetch failed for {container_path}: {err}"
    handle = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
    handle.write(base64.b64decode(encoded))
    handle.flush()
    handle.close()
    return Path(handle.name)


def _resolve_generated_pptx_for_validation(*, base_url: str, assistant_text: str) -> Path:
    raw_reference, pptx_path = _extract_generated_pptx_reference(assistant_text)
    if pptx_path.is_file():
        return pptx_path

    parsed = urlparse(base_url)
    hostname = (parsed.hostname or "").strip().lower()
    if raw_reference.startswith("/app/") and hostname not in {"", "127.0.0.1", "localhost"}:
        return _download_remote_container_file(base_url=base_url, container_path=Path(raw_reference))

    return pptx_path


@pytest.mark.e2e
@pytest.mark.integration
def test_live_agent_cmp_pending_returns_grounded_results() -> None:
    base_url, username, password = _require_live_e2e()
    session, token, session_key = _login_and_create_thread(
        base_url=base_url,
        username=username,
        password=password,
    )

    pending = _run_round_with_transient_retry(
        session,
        base_url=base_url,
        token=token,
        session_key=session_key,
        message="查下CMP 里目前所有待审批",
    )
    _assert_common_runtime_shape(pending)
    _assert_real_tool_call(pending, "smartcmp_list_pending")
    request_ids = _extract_cmp_request_ids(pending.assistant_text)
    assert request_ids, pending.assistant_text
    assert any(
        marker in pending.assistant_text
        for marker in ("审批", "待审批", "approv", "Workflow ID", "Request ID")
    ), pending.assistant_text

    print(f"LIVE_AGENT_TIMING scenario=cmp_pending elapsed={pending.wall_seconds}s")
    session.close()


@pytest.mark.e2e
@pytest.mark.integration
def test_live_agent_cmp_detail_returns_grounded_results() -> None:
    base_url, username, password = _require_live_e2e()
    session, token, session_key = _login_and_create_thread(
        base_url=base_url,
        username=username,
        password=password,
    )

    pending = _run_round_with_transient_retry(
        session,
        base_url=base_url,
        token=token,
        session_key=session_key,
        message="查下CMP 里目前所有待审批",
    )
    _assert_common_runtime_shape(pending)
    _assert_real_tool_call(pending, "smartcmp_list_pending")
    request_ids = _extract_cmp_request_ids(pending.assistant_text)
    assert request_ids, pending.assistant_text

    target_request_id = request_ids[0]

    detail = _run_round_with_transient_retry(
        session,
        base_url=base_url,
        token=token,
        session_key=session_key,
        message=f"我要看下{target_request_id}的详情",
    )
    _assert_common_runtime_shape(detail)
    _assert_real_tool_call(detail, "smartcmp_get_request_detail")
    assert target_request_id in detail.assistant_text
    assert (
        "approvalId" in detail.assistant_text
        or "requestId" in detail.assistant_text
        or "Approval ID" in detail.assistant_text
        or "Request ID" in detail.assistant_text
        or "Workflow ID" in detail.assistant_text
    )

    print(f"LIVE_AGENT_TIMING scenario=cmp_detail elapsed={detail.wall_seconds}s")
    session.close()


@pytest.mark.e2e
@pytest.mark.integration
def test_live_agent_cmp_services_returns_catalog_results() -> None:
    base_url, username, password = _require_live_e2e()
    session, token, session_key = _login_and_create_thread(
        base_url=base_url,
        username=username,
        password=password,
    )

    services = _run_round_with_transient_retry(
        session,
        base_url=base_url,
        token=token,
        session_key=session_key,
        message="还有查下CMP 里目前有的服务目录",
    )
    _assert_common_runtime_shape(services)
    _assert_real_tool_call(services, "smartcmp_list_services")
    assert "Incident Ticket" in services.assistant_text
    assert "Machine Service" in services.assistant_text
    assert "VPC Service" in services.assistant_text
    assert "Support Service" in services.assistant_text

    print(f"LIVE_AGENT_TIMING scenario=cmp_services elapsed={services.wall_seconds}s")
    session.close()


@pytest.mark.e2e
@pytest.mark.integration
def test_live_agent_weather_query_returns_grounded_weather() -> None:
    base_url, username, password = _require_live_e2e()
    session, token, session_key = _login_and_create_thread(
        base_url=base_url,
        username=username,
        password=password,
    )

    weather = _run_round_with_transient_retry(
        session,
        base_url=base_url,
        token=token,
        session_key=session_key,
        message="帮我查下上海天气",
    )
    _assert_common_runtime_shape(weather)
    _assert_real_tool_call(weather, "openmeteo_weather")
    assert "Weather for 上海" in weather.assistant_text or "Weather for Shanghai" in weather.assistant_text
    assert "Current:" in weather.assistant_text
    assert "°C" in weather.assistant_text

    print(f"LIVE_AGENT_TIMING scenario=weather elapsed={weather.wall_seconds}s")
    session.close()


@pytest.mark.e2e
@pytest.mark.integration
def test_live_agent_weather_follow_up_short_location_reuses_context() -> None:
    base_url, username, password = _require_live_e2e()
    session, token, session_key = _login_and_create_thread(
        base_url=base_url,
        username=username,
        password=password,
    )

    beijing = _run_round_with_transient_retry(
        session,
        base_url=base_url,
        token=token,
        session_key=session_key,
        message="明天北京天气呢",
    )
    _assert_common_runtime_shape(beijing)
    _assert_real_tool_call(beijing, "openmeteo_weather")
    assert re.search(r"(北京|北京市)", beijing.assistant_text)

    shanghai = _run_round_with_transient_retry(
        session,
        base_url=base_url,
        token=token,
        session_key=session_key,
        message="上海呢",
    )
    _assert_common_runtime_shape(shanghai)
    _assert_real_tool_call(shanghai, "openmeteo_weather")
    assert re.search(r"(上海|Shanghai)", shanghai.assistant_text)
    assert "Weather for" in shanghai.assistant_text or "天气" in shanghai.assistant_text

    print(
        "LIVE_AGENT_TIMING "
        f"scenario=weather_follow_up_short_location "
        f"beijing={beijing.wall_seconds}s "
        f"shanghai={shanghai.wall_seconds}s"
    )
    session.close()


@pytest.mark.e2e
@pytest.mark.integration
def test_live_agent_public_park_query_answers_directly() -> None:
    base_url, username, password = _require_live_e2e()
    session, token, session_key = _login_and_create_thread(
        base_url=base_url,
        username=username,
        password=password,
    )

    parks = _run_round_with_transient_retry(
        session,
        base_url=base_url,
        token=token,
        session_key=session_key,
        message="我想查下上海周边的骑行公园",
    )
    _assert_common_runtime_shape(parks, require_tool=False)
    assert not any(str(item).strip() for item in parks.errors)
    assert re.search(r"(崇明|青浦|浦东|松江|东平|青西|世纪公园)", parks.assistant_text)
    assert "骑行" in parks.assistant_text

    print(f"LIVE_AGENT_TIMING scenario=park_recommendation elapsed={parks.wall_seconds}s")
    session.close()


@pytest.mark.e2e
@pytest.mark.integration
def test_live_agent_cmp_competitors_query_answers_directly() -> None:
    base_url, username, password = _require_live_e2e()
    session, token, session_key = _login_and_create_thread(
        base_url=base_url,
        username=username,
        password=password,
    )

    competitors = _run_round_with_transient_retry(
        session,
        base_url=base_url,
        token=token,
        session_key=session_key,
        message="你查下有CMP有哪些类似的产品",
    )
    _assert_common_runtime_shape(competitors, require_tool=False)
    assert "cmp" in competitors.assistant_text.lower() or "云管理平台" in competitors.assistant_text
    assert re.search(
        r"(ServiceNow|VMware|CloudBolt|Morpheus|Flexera|Nutanix|BMC|IBM|博云|云霁|行云管家)",
        competitors.assistant_text,
        re.IGNORECASE,
    )

    print(f"LIVE_AGENT_TIMING scenario=cmp_competitors elapsed={competitors.wall_seconds}s")
    session.close()


@pytest.mark.e2e
@pytest.mark.integration
def test_live_agent_cmp_pending_then_ppt_creates_real_pptx() -> None:
    base_url, username, password = _require_live_e2e()
    session, token, session_key = _login_and_create_thread(
        base_url=base_url,
        username=username,
        password=password,
    )

    pending = _run_round_with_transient_retry(
        session,
        base_url=base_url,
        token=token,
        session_key=session_key,
        message="查一个 cmp 所有待审批的申请",
    )
    _assert_common_runtime_shape(pending)
    _assert_real_tool_call(pending, "smartcmp_list_pending")

    ppt = _run_round_with_transient_retry(
        session,
        base_url=base_url,
        token=token,
        session_key=session_key,
        message="将这些申请写入一个新的PPT",
    )
    _assert_common_runtime_shape(ppt)
    _assert_real_tool_call(ppt, "pptx_create_deck")
    assert ".pptx" in ppt.assistant_text.lower()
    pptx_path = _resolve_generated_pptx_for_validation(
        base_url=base_url,
        assistant_text=ppt.assistant_text,
    )
    assert pptx_path.is_file(), f"PPTX file not created: {pptx_path}"

    presentation = Presentation(str(pptx_path))
    assert len(presentation.slides) >= 3
    title_texts = [shape.text for shape in presentation.slides[0].shapes if hasattr(shape, "text")]
    assert any(re.search(r"(CMP|待审批|PPT)", text) for text in title_texts)

    print(
        "LIVE_AGENT_TIMING "
        f"scenario=cmp_pending_then_ppt "
        f"pending={pending.wall_seconds}s "
        f"ppt={ppt.wall_seconds}s"
    )
    session.close()


@pytest.mark.e2e
@pytest.mark.integration
def test_live_agent_cmp_pending_then_english_ppt_follow_up_creates_real_pptx() -> None:
    base_url, username, password = _require_live_e2e()
    session, token, session_key = _login_and_create_thread(
        base_url=base_url,
        username=username,
        password=password,
    )

    pending = _run_round_with_transient_retry(
        session,
        base_url=base_url,
        token=token,
        session_key=session_key,
        message="查下CMP现在的审批数据",
    )
    _assert_common_runtime_shape(pending)
    _assert_real_tool_call(pending, "smartcmp_list_pending")

    ppt = _run_round_with_transient_retry(
        session,
        base_url=base_url,
        token=token,
        session_key=session_key,
        message="write the request data into a PPT",
    )
    _assert_common_runtime_shape(ppt)
    _assert_real_tool_call(ppt, "pptx_create_deck")
    assert ".pptx" in ppt.assistant_text.lower()
    assert ".txt" not in ppt.assistant_text.lower()
    pptx_path = _resolve_generated_pptx_for_validation(
        base_url=base_url,
        assistant_text=ppt.assistant_text,
    )
    assert pptx_path.is_file(), f"PPTX file not created: {pptx_path}"

    presentation = Presentation(str(pptx_path))
    assert len(presentation.slides) >= 3

    print(
        "LIVE_AGENT_TIMING "
        f"scenario=cmp_pending_then_english_ppt "
        f"pending={pending.wall_seconds}s "
        f"ppt={ppt.wall_seconds}s"
    )
    session.close()
